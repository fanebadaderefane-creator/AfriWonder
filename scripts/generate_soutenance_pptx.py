from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent.parent / "docs" / "soutenance" / "Soutenance_.pptx"


def style_title(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 34, 64)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Soutenance de Projet : AfriWonder Super-App"
    style_title(slide.shapes.title)

    subtitle = (
        "Projet AfriWonder — Écosystème numérique multi-services\n\n"
        "Candidat : FANE ABDOULAYE (Élève Ingénieur)\n"
        "Encadrant : Pr. HAMZA KHALFI\n"
        "Filière : Management et Gouvernance des Systèmes d’Information (MGSI)\n"
        "Établissement : École Nationale des Sciences Appliquées de Khouribga (ENSAK)"
    )
    slide.placeholders[1].text = subtitle
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.size = Pt(18)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6), Inches(13.33), Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(16, 34, 64)
    band.line.fill.background()
    band.text_frame.text = "Architecture | Gouvernance SI | Performance | Sécurité | Industrialisation"
    for p in band.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    if note:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(5.8), Inches(4.9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(232, 240, 255)
        box.line.color.rgb = RGBColor(100, 130, 190)
        box.text_frame.text = note
        for p in box.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(20, 40, 80)


def add_architecture_illustration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Architecture Logique du Système (Illustration)"

    def node(x, y, w, h, text, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.color.rgb = RGBColor(60, 60, 60)
        s.text_frame.text = text
        for p in s.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        return s

    node(0.5, 1.4, 2.7, 1.0, "PWA React/Vite", RGBColor(65, 105, 225))
    node(0.5, 3.0, 2.7, 1.0, "Mobile Expo", RGBColor(65, 105, 225))
    node(4.0, 2.2, 2.8, 1.2, "API Express\n(TypeScript)", RGBColor(0, 128, 128))
    node(7.7, 1.4, 2.5, 1.0, "Prisma ORM", RGBColor(123, 63, 0))
    node(10.5, 1.4, 2.3, 1.0, "PostgreSQL", RGBColor(47, 79, 79))
    node(7.7, 3.0, 2.5, 1.0, "Redis / Socket", RGBColor(178, 34, 34))
    node(10.5, 3.0, 2.3, 1.0, "S3/R2 Media", RGBColor(178, 34, 34))

    for (x, y, w, h) in [(3.2, 1.85, 0.8, 0.08), (3.2, 3.45, 0.8, 0.08), (6.8, 2.75, 0.9, 0.08), (10.2, 1.85, 0.3, 0.08), (6.8, 3.35, 0.9, 0.08), (10.2, 3.35, 0.3, 0.08)]:
        c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(60, 60, 60)
        c.line.fill.background()

    legend = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.4))
    legend.text_frame.text = "Séparation en couches : Client -> API -> Services -> Données.\nFlux temps réel via Socket.IO + Redis adapter (multi-instance)."
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(16)


def add_payment_sequence_illustration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Séquence Paiement (Illustration)"

    cols = ["Client", "API Payments", "Stripe/Orange", "Wallet"]
    xs = [0.6, 3.6, 6.7, 9.9]
    for label, x in zip(cols, xs):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.1), Inches(2.3), Inches(0.7))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(16, 34, 64)
        b.text_frame.text = label
        for p in b.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(13)

    steps = [
        (0.9, 2.1, 5.8, "1) Init paiement"),
        (3.9, 2.7, 8.8, "2) Création intent / transaction"),
        (7.0, 3.3, 4.6, "3) Webhook de confirmation"),
        (4.0, 3.9, 11.5, "4) MAJ Wallet + statut commande"),
    ]

    for x1, y, x2, txt in steps:
        w = x2 - x1
        line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y), Inches(w), Inches(0.35))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(70, 130, 180)
        line.line.fill.background()
        t = slide.shapes.add_textbox(Inches(x1), Inches(y - 0.32), Inches(w), Inches(0.25))
        t.text_frame.text = txt
        t.text_frame.paragraphs[0].font.size = Pt(11)


def add_cicd_illustration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Pipeline CI/CD (Illustration)"

    stages = [
        "Push / PR",
        "Lint + Typecheck",
        "Tests Jest/Vitest",
        "Playwright E2E",
        "Build Docker",
        "Déploiement",
    ]
    colors = [
        RGBColor(44, 62, 80), RGBColor(52, 152, 219), RGBColor(46, 204, 113),
        RGBColor(241, 196, 15), RGBColor(230, 126, 34), RGBColor(155, 89, 182)
    ]
    x = 0.5
    for i, s in enumerate(stages):
        box = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(2.4), Inches(2.0), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        box.text_frame.text = s
        for p in box.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(12)
        x += 2.0

    txt = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(12.0), Inches(1.6))
    txt.text_frame.text = (
        "Contrôles qualité automatisés :\n"
        "- Budget taille PR\n"
        "- Lint + Typecheck bloquants\n"
        "- Tests unitaires/intégration + E2E"
    )
    for p in txt.text_frame.paragraphs:
        p.font.size = Pt(16)


def build_presentation():
    prs = Presentation()

    add_cover(prs)

    add_bullets(prs, "Sommaire de la Présentation", [
        "Introduction et Contexte du Projet",
        "Analyse Fonctionnelle et Métiers",
        "Architecture Système et Stack Technologique",
        "Conception de la Base de Données",
        "Développement Mobile et PWA",
        "Sécurité, Performance et DevOps",
        "Conclusion et Perspectives",
    ])

    add_bullets(prs, "Introduction au Projet AfriWonder", [
        "Genèse du projet : Répondre aux défis de la fragmentation numérique en Afrique",
        "Vision : Créer une \"Super-App\" unifiée (Social, Commerce, Finance, Services)",
        "Ambition : Centraliser l'expérience utilisateur dans un écosystème régional cohérent",
        "Valeur ajoutée : Accessibilité, monétisation locale et services de proximité",
    ], "Illustration suggérée : carte du parcours utilisateur unifié.")

    add_bullets(prs, "Problématique et Enjeux Stratégiques", [
        "Fragmentation des services mobiles en Afrique",
        "Difficultés d'accès aux infrastructures de paiement internationales",
        "Contraintes de connectivité et de performance des terminaux",
        "Nécessité d'une gouvernance centralisée des données (MGSI)",
    ])

    add_bullets(prs, "Analyse des Besoins Fonctionnels", [
        "Modules principaux : Feed vidéo, Marketplace, Messagerie temps réel",
        "Services intégrés : Wallet financier, Gestion de transports, Santé, Tontines",
        "Gestion des profils : Créateurs, Vendeurs, Prestataires, Administrateurs",
        "Fonctions transverses : Traduction, Chatbots, Notifications push",
    ], "Détail : couvre social, commerce et services transactionnels.")

    add_bullets(prs, "Cartographie des Cas d’Utilisation", [
        "Publication et interaction avec le contenu vidéo (Social)",
        "Cycle de vente complet : de l'annonce au paiement sécurisé (Marketplace)",
        "Gestion de portefeuille et transactions financières (Wallet)",
        "Modération et administration globale de la plateforme",
    ])

    add_bullets(prs, "Architecture Logique du Système", [
        "Modèle Client-Serveur distribué",
        "Communication via API REST sécurisée et WebSockets",
        "Découpage modulaire : Séparation des domaines métier au sein du backend",
        "Architecture en couches : Routes, Services, Accès aux données",
    ])
    add_architecture_illustration(prs)

    add_bullets(prs, "Stack Technologique : Le Choix de l'Efficacité", [
        "Langage unique : TypeScript pour le Frontend et le Backend",
        "Backend : Node.js avec le framework Express",
        "Frontend Web : React 18 avec l'outil de build Vite",
        "Mobile : React Native via le framework Expo (SDK 54+)",
    ])

    add_bullets(prs, "Architecture Backend et API", [
        "Gestion des routes structurée par domaine métier",
        "Validation rigoureuse des données avec la bibliothèque Zod",
        "Middleware de protection : Auth JWT, Rate Limiting, Helmet, CORS",
        "Gestion des fichiers et médias avec Multer et Sharp",
    ], "Explication : architecture orientée modules pour maintenir un périmètre large.")

    add_bullets(prs, "Modélisation des Données avec Prisma", [
        "Utilisation de Prisma ORM pour une productivité accrue",
        "Schéma relationnel PostgreSQL robuste (plus de 100 modèles)",
        "Gestion des migrations automatisée et typage fort des entités",
        "Optimisation des requêtes pour un volume de données massif",
    ])

    add_bullets(prs, "Structure du Schéma de Base de Données", [
        "L'entité User : Hub central reliant Vidéos, Wallet et Commandes",
        "Modélisation de la Marketplace : Produits, Catégories, Paniers et État des stocks",
        "Système de Wallet : Transactions, Soldes et Historique financier",
        "Gestion sociale : Abonnements, Commentaires, Likes et Stories",
    ])

    add_bullets(prs, "Stratégie Frontend : PWA vs Mobile Natif", [
        "PWA (Vite/React) : Accessibilité web rapide et légèreté (SEO, Reach)",
        "Application Mobile (Expo) : Expérience native, Caméra, Notifications, Offline",
        "Partage de la logique métier et des types entre les deux clients",
        "Design System unifié avec Tailwind CSS et composants réutilisables",
    ])

    add_bullets(prs, "Gestion de l'État et Récupération des Données", [
        "TanStack Query (React Query) pour la gestion du cache serveur",
        "Zustand pour l'état global léger (Authentification, Thème)",
        "Persistance des données locales pour une expérience fluide",
        "Optimistic UI pour les interactions sociales (Likes, Commentaires)",
    ])

    add_bullets(prs, "Communication en Temps Réel", [
        "Intégration de Socket.IO pour la messagerie instantanée",
        "Gestion des salons de discussion et présence en ligne",
        "Notifications en temps réel pour les activités sociales et commerciales",
        "Architecture scalable avec adaptateur Redis pour le multi-instance",
    ])

    add_bullets(prs, "Intégration des Services de Paiement", [
        "Passerelle multi-paiements : Stripe et solutions locales (Orange Money)",
        "Gestion sécurisée des Webhooks pour la validation des transactions",
        "Workflow de paiement : Initialisation, Confirmation, Mise à jour du Wallet",
        "Sécurisation des transactions et traçabilité financière",
    ])
    add_payment_sequence_illustration(prs)

    add_bullets(prs, "Services de Traduction et Intelligence Artificielle", [
        "API de traduction intégrée (LibreTranslate / MyMemory)",
        "Système de Chatbots pour l'assistance automatisée",
        "Pipeline de génération de sous-titres (STT) pour les vidéos",
        "Modération automatisée basée sur des règles et listes de contrôle",
    ], "Note : certaines briques IA sont en mode progressif / placeholder dans le code.")

    add_bullets(prs, "Sécurité du Système d'Information", [
        "Authentification par Tokens JWT (Access & Refresh Tokens)",
        "Protection contre les attaques courantes (XSS, CSRF, Injection SQL via Prisma)",
        "Proxy média sécurisé avec liste blanche de domaines",
        "Gestion stricte des variables d'environnement et secrets",
    ])

    add_bullets(prs, "Performance et Optimisations Techniques", [
        "Code Splitting et Lazy Loading côté Frontend pour réduire le temps de chargement",
        "Compression Gzip/Brotli et mise en cache Redis côté Backend",
        "Optimisation des images et streaming vidéo par fragments",
        "Surveillance des performances avec Sentry et métriques Prometheus",
    ])

    add_bullets(prs, "Stratégie DevOps et Déploiement", [
        "Conteneurisation complète avec Docker et Docker Compose",
        "Gestion de 3 réplicas pour le Backend afin d'assurer la haute disponibilité",
        "Serveur Nginx utilisé comme Reverse Proxy et gestionnaire TLS (Certbot)",
        "Base de données PostgreSQL 15 et cache Redis 7",
    ])

    add_bullets(prs, "Automatisation et Qualité du Code (CI/CD)", [
        "GitHub Actions pour l'intégration continue (CI)",
        "Vérification automatique de la qualité du code (Lint, Typecheck)",
        "Tests unitaires et d'intégration avec Jest et Vitest",
        "Contraintes de taille de Pull Request pour maintenir la maintenabilité",
    ])
    add_cicd_illustration(prs)

    add_bullets(prs, "Défis Techniques et Solutions Apportées", [
        "Gestion de la complexité d'un schéma Prisma massif",
        "Optimisation du lecteur vidéo sous contraintes réseau instables",
        "Synchronisation des états d'authentification entre PWA et Mobile",
        "Sécurisation du proxy média pour éviter les failles SSRF",
    ])

    add_bullets(prs, "Limites Actuelles du Projet", [
        "Monolithe applicatif pouvant présenter des défis de scalabilité verticale",
        "Dépendance forte à certains services tiers pour l'IA et la traduction",
        "Parties placeholders pour les fonctions avancées de STT (Speech-to-Text)",
        "Volume important du schéma DB ralentissant l'onboarding développeur",
    ])

    add_bullets(prs, "Améliorations Futures et Évolutions", [
        "Transition vers une architecture microservices pour les domaines critiques",
        "Finalisation de l'intégration native des modèles d'IA pour le STT/NLP",
        "Renforcement des tests de charge et de montée en charge (Stress Testing)",
        "Développement d'un module de \"Feature Flags\" pour le déploiement progressif",
    ])

    add_bullets(prs, "Compétences Acquises et Bilan Personnel", [
        "Maîtrise du cycle de vie complet d'une application Full-Stack",
        "Expertise en gouvernance des systèmes d’information complexes",
        "Capacité à concevoir des architectures scalables et sécurisées",
        "Gestion de projet technique sous standards de qualité industrielle",
    ])

    add_bullets(prs, "Conclusion Générale", [
        "AfriWonder : Une réponse technologique concrète aux besoins du marché africain",
        "Un projet d'ingénierie complet alliant innovation et rigueur technique",
        "Validation des acquis de la filière MGSI de l'ENSAK",
        "Prêt pour une phase d'industrialisation et de mise sur le marché",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build_presentation()
