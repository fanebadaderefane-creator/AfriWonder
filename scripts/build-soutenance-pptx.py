# -*- coding: utf-8 -*-
"""Génère docs/soutenance/Soutenance_.pptx à partir du contenu de soutenance (analyse repo)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "soutenance" / "Soutenance_.pptx"


def add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass


def add_bullets(prs, title: str, bullets: list[str], font_pt: float = 14):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)
    return slide


def add_two_column_bullets(prs, title: str, left_title: str, left: list, right_title: str, right: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.75))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True

    lbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.4), Inches(5.5))
    ltf = lbox.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.size = Pt(13)
    for b in left:
        p = ltf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(11)

    rbox = slide.shapes.add_textbox(Inches(5.1), Inches(1.15), Inches(4.4), Inches(5.5))
    rtf = rbox.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.size = Pt(13)
    for b in right:
        p = rtf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(11)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 0 — couverture
    add_title_slide(
        prs,
        "AfriWonder — Soutenance de fin d'études",
        "Étudiant : Abdoulaye Fanel\nEncadrant : Professeur Hamza Kalfi\nAnalyse technique du dépôt (PWA + API + mobile Expo)",
    )

    # --- 1. PRÉSENTATION GÉNÉRALE ---
    add_bullets(
        prs,
        "1. Présentation générale — Projet et objectif",
        [
            "Nom : AfriWonder — super-application sociale et commerciale orientée Afrique de l'Ouest.",
            "Objectif principal : plateforme unifiée (vidéo type feed, messaging, marketplace, services, paiements, lives…).",
            "Problématique : fragmenter moins l'expérience utilisateur (réseau, commerce, services) sur des connexions contraintes.",
            "Contexte technique : monorepo avec PWA React (`src/`), backend Express (`backend/src/`), mobile Expo RN (`frontend/`).",
            "Public cible : utilisateurs mobiles (prioritaire), marchés type Mali puis extension continentale.",
            "Références dépôt : `AGENTS.md`, `docs/ARCHITECTURE.md`, `backend/package.json`, `frontend/package.json`, racine `package.json`.",
        ],
    )
    add_bullets(
        prs,
        "1bis. Valeur ajoutée et différenciation",
        [
            "Un seul écosystème API central (`backend`) consommé par plusieurs clients.",
            "Pensée « réseau faible » : timeouts longs mobile (`frontend/src/api/client.ts`), mode économiseur données (`frontend/src/dataSaver/`).",
            "Couverture fonctionnelle très large (schéma Prisma volumineux) : modularité métier forte côté API.",
            "Différenciation : combinaison TikTok-like + marketplace + verticals locaux — à valider fonctionnellement écran par écran.",
        ],
    )

    # --- 2. ANALYSE FONCTIONNELLE ---
    add_bullets(
        prs,
        "2. Analyse fonctionnelle — périmètre observé dans le code",
        [
            "Principaux domaines montés sur l'API Express : voir liste d'imports `backend/src/app.ts` (videos, marketplace, paiements, live, messaging, géolocalisation, etc.).",
            "Messagerie : routes `messages.routes.ts`, groupes dans `messageGroup.service.ts`.",
            "Feed / annonces / feed dédiés : `feed.routes.ts`, `ads.routes.ts`.",
            "Comptes et rôles : champ `User.role` (Prisma), routes admin et modération dédiées.",
            "Chatbot : entités `ChatBot` exposées en lecture via `chatbot.routes.ts` (pas de LLM génératif dans ces routes — catalogue en base).",
            "⚠ Une superficie très large impose des validations terrain : tout endpoint n'équivaut pas à une feature UX « terminée » côté client.",
        ],
    )
    add_bullets(
        prs,
        "2bis. Permissions, erreurs et cas limites",
        [
            "Authentification : middleware `authenticate` / `optionalAuth` (`middleware/auth.js`) utilisé selon routes.",
            "Validation : helpers Zod réutilisables (`utils/zodValidation.js`, schémas dédiés).",
            "Gestion erreurs : chaîne middleware Express jusqu'à handlers (cf. conventions backend).",
            "Rate limiting : `middleware/rateLimiting.ts` — variantes général/auth/paiement/upload/admin/webhook.",
            "Cas limite réseau mobile : résolution dynamique backend dev Android (`frontend/src/config/backendBase.ts`).",
        ],
    )

    # --- 3. ARCHITECTURE ---
    add_bullets(
        prs,
        "3. Architecture globale — couches et communication",
        [
            "Client web PWA : Vite + React 18 dans `src/` — build `vite.config.js`, routing `react-router-dom`.",
            "Client mobile : Expo Router `frontend/app/` + React Native 0.81 + React 19 (cf. `frontend/package.json`).",
            "API : Node 20+ + Express 4 dans `backend/src/index.ts` / `app.ts`.",
            "Données : PostgreSQL via Prisma ORM (`backend/prisma/schema.prisma`).",
            "Temps réel : Socket.IO client (PWA + mobile deps) ; serveur `socket.io` + adaptateur Redis optionnel backend.",
            "Fichiers : AWS SDK S3 + presign (`backend/package.json`) ; CDN évoqué dans Vite (`VITE_MEDIA_PRECONNECT`).",
        ],
    )
    add_bullets(
        prs,
        "3bis. Choix architecture — avantages / limites / alternatives",
        [
            "Choix monolithique Express : simplicité de déploiement, introspection Swagger `/api-docs`, mutualisation middleware sécurité.",
            "Alternative microservices : meilleure isolation par domaine, mais coût ops et complexité synchrone/async — non retenue comme base.",
            "Séparation PWA/mobile : optimise UX natives (notifications, IAP, caméra, Agora RN) tout en gardant un contrat REST unique.",
            "Dette observable : dossier TypeScript `functions/` (28 fichiers) en parallèle — non monté depuis `backend/src` (pas d'import repéré) : peut servir historique/serverless hors chemin critique actuel.",
        ],
    )

    # --- 4. FRONTEND WEB (PWA) ---
    add_bullets(
        prs,
        "4. Analyse frontend PWA (`src/` à la racine)",
        [
            "Framework : React — pas Next.js (build Vite SPA + PWA `vite-plugin-pwa`). Pourquoi : build rapide, contrôle précis bundles, PWAs hors SSR.",
            "Données : TanStack Query + persistance IndexedDB/sync storage deps racine.",
            "UI : Tailwind CSS + grande surface Radix UI + Framer Motion + composants métier vidéo (`src/components/video/VideoCard.jsx`).",
            "Médias : Agora Web SDK (`agora-rtc-sdk-ng`), HLS (`hls.js`), streaming temps réel pour lives.",
            "Sécu client : sanitize HTML potentiel avec `dompurify` dependency.",
            "Observabilité : `@sentry/react`, analytics PostHog (`posthog-js`).",
            "Tests & qualité : Vitest, Playwright (`tests/e2e/`), Lighthouse CI scripts racine.",
        ],
    )

    # --- 4b MOBILE ---
    add_bullets(
        prs,
        "4bis. Analyse frontend mobile Expo (`frontend/`)",
        [
            "Navigation : expo-router (`app/(tabs)/…`) + React Navigation packages.",
            "État : Zustand + TanStack Query ; stockage sensible `expo-secure-store` (cf. `secureStorage`).",
            "Média : expo-video, vision-camera, Agora (`react-native-agora`), WebRTC.",
            "Paiements : `react-native-iap` (stores) + appels API proxy alignés PWA.",
            "Qualité : `npm run verify` = lint Expo + typecheck + Vitest couverture (seuil défini `frontend/vitest.config.ts`).",
        ],
    )

    # --- 5. BACKEND ---
    add_bullets(
        prs,
        "5. Analyse backend Express (`backend/src/`)",
        [
            "Stack : TypeScript, Express, Prisma 7, engines Node >= 20.19.",
            "Sécurité HTTP : `helmet`, `cors`, `cookie-parser`, limites `express-rate-limit` avec store Redis optionnel.",
            "Protection requêtes : `sanitizeInputMiddleware`, `csrfProtectionMiddleware`, `cachePolicyMiddleware` (`requestProtection.middleware.ts`).",
            "Observabilité : Sentry Node/Profiling, métriques Prometheus exposition via service dédié, timeouts API (`observability.middleware.ts`).",
            "Docs API : Swagger UI + `swagger-jsdoc` — vérité fonctionnelle `GET /api/openapi.json` (cf. ARCHITECTURE).",
            "Paiements : Stripe + webhooks dédiés, Orange Money, PDF `pdfkit`, qrcode.",
            "Push : web-push dépendances ; Firebase Admin pour notifications mobiles probables.",
        ],
    )

    # --- 6. BASE DE DONNÉES ---
    add_bullets(
        prs,
        "6. Base de données — PostgreSQL + Prisma",
        [
            "SGBDR relationnel PostgreSQL → intégrité transactionnelle, JOINs riches, migrations versionnées Prisma.",
            "Schéma massif (>6000 lignes) : forte modularité fonctionnelle (users, commerce, géo, santé, events, crowdfunding…).",
            "Relations : très nombreuses FK depuis `User` — modèle super-app (« many concerns one user core »).",
            "Pourquoi SQL vs document : reporting, agrégats, workflows transactionnels paiement/commande.",
            "Indexation / perf : présence probable d'indexes sur champs `@unique`/FK — optimisation continue à profiler sur charges réelles.",
            "⚠ Une surface aussi large impose rigueur migrations + revues performance (indexes ciblés, N+1).",
        ],
    )

    # --- 7. IA ---
    add_bullets(
        prs,
        "7. IA / ML — usages réels identifiés dans le code",
        [
            "Transcription vocale serveur : OpenAI Whisper (`whisper-1` ou `OPENAI_TRANSCRIPTION_MODEL`), via `backend/src/utils/whisperTranscription.ts`.",
            "Intégration messages 1–1 et groupes : `message.service.ts`, `messageGroup.service.ts`, routes transcription dédiées.",
            "Live STT hôte : `liveStt.service.ts` même API Whisper avec clés serveur uniquement.",
            "Traduction `/api`-translate : LibreTranslate (+ fallback MyMemory) — pas de LLM génératif propriétaire dans `translate.routes.ts`.",
            "Sous-titres vidéo : `subtitle.service.ts` prévoit STT externe — placeholder/documenté dans le fichier.",
            "Pas de pipeline RAG / embeddings repéré dans backend principal — ne pas attribuer hors preuve fichier.",
            "⚠ Risques IA : disponibilité API, confidentialité audio, facturation tokens — hors scope si clé absente (.env documenté).",
        ],
    )

    # --- 8. DEVOPS ---
    add_bullets(
        prs,
        "8. DevOps & déploiement",
        [
            "Docker : `docker-compose.prod.yml` — services backend (réplicas swarm), frontend Nginx TLS, Postgres 15, Redis 7, certbot renewal.",
            "Fichiers complémentaires : `docker-compose.replication.yml`, `docker-compose.scaling-1m.yml` (capacité / charge).",
            "CI/CD : `.github/workflows/ci.yml` — gates lint/typecheck, tests backend Jest, Vitest mobile, Playwright e2e, budgets PR (≤400 lignes).",
            "Variables : `.env` multiples (backend/frontend) — secrets jamais commit (cf. règles sécurité dépôt).",
            "Git : stratégie branches `main`/`develop` + PR (AGENTS.md).",
            "Monitoring : Sentry (apps), scripts load-test backend (`load-test-node.js`).",
        ],
    )

    # --- 9. SÉCURITÉ ---
    add_bullets(
        prs,
        "9. Sécurité — mesures observées",
        [
            "Auth : JWT (`jsonwebtoken` + `jose` dépendances) — refresh token côté produit à vérifier dans `auth` services/routes.",
            "Mots de passe : bcryptjs (selon deps backend — aligné politique sécurité interne).",
            "Transport : HTTPS en prod via Nginx + Let's Encrypt template (`docker-compose.prod.yml`).",
            "Entrées : Zod validation, sanitization middleware, anti-bot middleware dédié.",
            "Protection abus : rate-limit par classe d'endpoint ; webhooks exclus des limites générales.",
            "XSS / CSRF : layer `requestProtection.middleware` + pratiques React (éviter dangerouslySetInnerHTML non contrôlé).",
            "Stockage tokens mobile : SecureStore (Expo).",
        ],
    )

    # --- 10. PERF ---
    add_bullets(
        prs,
        "10. Performance & optimisation",
        [
            "Frontend PWA : code splitting Vite, compression plugin optionnelle, préconnect API/CDN (`vite.config.js`).",
            "Cache HTTP : middleware politique cache côté API.",
            "Images : Sharp backend + expo-image mobile (deps).",
            "Listes mobiles : `@shopify/flash-list` pour virtualisation.",
            "Requêtes : TanStack Query (stale-time, dedup), stratégie cache feed documentée dans règles internes `.cursor/rules/`.",
            "DB : pagination attendue sur endpoints listing — à valider endpoint par endpoint.",
        ],
    )

    add_bullets(
        prs,
        "10bis. Parcours utilisateur (exemples techniques)",
        [
            "Onboarding / auth : routes `auth.routes.ts` + stockage token côté clients (SecureStore mobile, persistance web selon implémentation).",
            "Feed vidéo PWA : `src/pages/Home.jsx` + `VideoCard.jsx` — interactions like/scroll documentées (règles internes feed).",
            "Création contenu mobile : écrans `frontend/app/` (upload) + timeouts multipart documentés dans `apiClient`.",
            "Marketplace : routes `products`, `orders`, `cart`, `payments` montées dans `app.ts`.",
            "Support & légal : `legal.routes.ts`, `publicPages.routes.ts` (pages statiques réglementaires Play Store évoquées en commentaire).",
        ],
    )

    add_bullets(
        prs,
        "10ter. État global & données (comparaison des choix)",
        [
            "TanStack Query (web + mobile) : cache serveur, invalidation, requêtes dédupliquées — adapté aux APIs REST et au offline partiel PWA.",
            "Zustand (mobile) : état UI/léger global sans boilerplate Redux — complète React Query pour l’état non serveur.",
            "Pourquoi pas Redux Toolkit partout : complexité reduite pour prototypage rapide ; trade-off : discipline pour ne pas disperser l’état mutable.",
            "Client API unifié côté web : `src/api/expressClient.js` (init dans `src/main.jsx`) vs mobile `frontend/src/api/client.ts` sur `/api/proxy`.",
        ],
    )

    add_bullets(
        prs,
        "10quater. Formulaires, validation, accessibilité (PWA)",
        [
            "Formulaires : `react-hook-form` + `@hookform/resolvers` + `zod` (dépendances racine `package.json`).",
            "Validation : schémas Zod partagés côté mental model client/serveur (serveur = source de vérité dans `backend`).",
            "Accessibilité : Radix UI primitives (comportement clavier, focus) — réduit le coût d’implémentation a11y de base.",
            "ErrorBoundary chargement appli : `src/main.jsx` référence `ErrorBoundary` pour limiter crash total.",
        ],
    )

    # --- 11. DIFFICULTÉS ---
    add_bullets(
        prs,
        "11. Difficultés techniques (inférées du code & commentaires)",
        [
            "Lecteur vidéo web : stabilisation buffer / autoplay documentée comme zone sensible (`VideoCard.jsx` + règles Cursor).",
            "Réseau mobile dev : bugs passés MEmu/Metro — logique `probeAndroidDevBackendOrigin` (`backendBase.ts`).",
            "Interop CJS/ESM : contournements explicites Recharts + Vite (`vite.config.js`).",
            "Surface API vs clients : parité mobile/PWA exigée par la gouvernance (`mobile-frontend-quality-api-parity.mdc`).",
            "Complexité Prisma : migrations et perfs requissent discipline (risque N+1 sur routes non auditées).",
        ],
    )

    # --- 12-14 ---
    add_bullets(
        prs,
        "12. Compétences acquises",
        [
            "Full-stack TypeScript/JS : React, React Native, Node, Prisma.",
            "Architecture produit : multi-clients, proxy API unifié (`/api/proxy` côté mobile).",
            "Qualité : tests automatisés (Jest, Vitest, Playwright), CI GitHub Actions.",
            "DevOps : Docker compose prod, Redis, Postgres, TLS.",
            "Sécurité : auth, rate limit, validation, observabilité.",
            "Domaine média : HLS, Agora, uploads S3-presign.",
        ],
    )
    add_bullets(
        prs,
        "13. Limites actuelles (honnêteté technique)",
        [
            "Superficie fonctionnelle énorme — risque de features partielles ou placeholders (ex. parties STT/sous-titres).",
            "Héritage `functions/` non câblé au runtime Express central — clarification documentation produit nécessaire.",
            "Dépendance services externes (Whisper, Stripe, Agora, FCM) — disponibilité et conformité données.",
            "Tests e2e ne couvrent pas exhaustivement tous les modules verticaux.",
            "Dette : fichiers longs possibles (standards internes 300 lignes cible).",
        ],
    )
    add_bullets(
        prs,
        "14. Améliorations futures réalistes",
        [
            "Cartographier maturité feature par module (score + owners).",
            "Durcir parcours critique (auth, wallet, feed) avec SLO / alertes.",
            "Étendre tests contrat OpenAPI ↔ clients SDK généré.",
            "Optimiser indexes SQL sur requêtes chaudes (APM + EXPLAIN).",
            "Évaluer extraction progressive de domaines (ex. paiements) si charge ops le justifie.",
            "Finaliser stratégie E2E messagerie (commentaires Prisma `messaging_e2e_enabled`).",
        ],
    )

    # --- Diagrammes (description) ---
    add_bullets(
        prs,
        "15. Diagrammes recommandés pour l'oral",
        [
            "Architecture globale : PWA + Expo → Express → PostgreSQL / Redis / S3 / Agora (cf. mermaid `docs/ARCHITECTURE.md`).",
            "Flux utilisateur : inscription → JWT → feed TanStack Query → interactions (like, commentaire, upload).",
            "Séquence API : refresh token / upload multipart (commentaires `apiClient` mobile).",
            "Modèle données : extrait simplifié User + Video + Order (ne pas projector tout Prisma en soutenance).",
            "Déploiement : schéma Nginx + backend répliqué + Postgres + Redis (`docker-compose.prod.yml`).",
        ],
    )

    add_bullets(
        prs,
        "16. Tests, qualité logicielle & outillage",
        [
            "Backend : Jest + Supertest, scripts `test`, `test:smoke`, `test:coverage` (`backend/package.json`).",
            "PWA : Vitest + Testing Library (`package.json` racine).",
            "Mobile : Vitest avec couverture gate CI (`frontend/vitest.config.ts`).",
            "E2E web : Playwright — jeux de tests ciblés dans `test:e2e:ci` (feed, navigation, paiement, Firefox layout).",
            "Mobile E2E : Maestro YAML (`frontend/maestro/`).",
            "Qualité transverse : ESLint backend/frontend/racine, typecheck TS, scripts `verify:quality-gates`, `verify:delivery`.",
        ],
    )

    add_bullets(
        prs,
        "17. Intégrations & services tiers (inventaire dépendances)",
        [
            "Paiements : Stripe, Orange Money (webhooks testés dans `backend/src/payments/__tests__/`).",
            "Vidéo / live : Agora (`agora-token` serveur, SDK web & RN client).",
            "Stockage objet : `@aws-sdk/client-s3`, presigner uploads.",
            "Auth supplémentaire : Supabase client (dépôt), Firebase Admin (push), Sign in Apple (Expo module).",
            "Cartes & carto : `react-native-maps`, `react-leaflet` (PWA).",
            "Observabilité : Sentry (Node, React, RN).",
            "Analytics : PostHog (initialisation lazy `src/main.jsx`).",
        ],
    )

    add_bullets(
        prs,
        "18. Conformité produit, légal & gouvernance du code",
        [
            "Pages légales & confidentialité : routes dédiées + exigence stores mentionnée (`publicPages.routes.ts`).",
            "Modération : `moderation.routes.ts` et signalements utilisateurs dans le schéma Prisma.",
            "Privacy : `privacy.routes.ts` dédiées.",
            "Gouvernance : `AGENTS.md` + scripts d’audit (`verify-engineering-standards`, taille PR CI).",
            "Documentation longue durée : `docs/ENGINEERING_STANDARDS.md`, `SECURITY.md`, runbooks.",
        ],
    )

    # --- 19. Code quality ---
    add_bullets(
        prs,
        "19. Analyse critique du code (dépôt)",
        [
            "Points forts : conventions AGENTS.md, gates CI, tests multi-couches, middleware sécurité structuré.",
            "Modularité backend : découpage `routes/`, `services/`, `middleware/`.",
            "Risque : cohérence long terme entre surface Prisma et UX réelle — discipline produit requise.",
            "Duplication potentielle : clients multiples (PWA/mobile) — centraliser client API généré pourrait réduire écarts.",
            "Maintenabilité : documentation vivante (`docs/`) et scripts `verify:*` racine.",
        ],
    )

    # --- 20. Conclusion ---
    add_bullets(
        prs,
        "20. Conclusion professionnelle",
        [
            "Bilan : plateforme ambitieuse combinant réseau social vidéo, commerce et services verticaux sur stack moderne cohérente.",
            "Résultats techniques : API Express riche, clients PWA + Expo, persistance PostgreSQL robuste, chaîne qualité CI.",
            "Compétences : conception full-stack, sécurité, perf réseaux contraints, intégration paiements/temps réel.",
            "Perspectives : industrialiser la carte des features, mesurer perf réelle terrain, poursuivre durcissement sécurité & conformité données.",
            "⚠ Synthèse basée exclusivement sur l'analyse du dépôt local — comportement prod dépend aussi du déploiement et secrets.",
        ],
    )

    # Annexes citations fichiers clés
    add_two_column_bullets(
        prs,
        "Annexe — Fichiers à citer au jury",
        "Backend / Infra",
        [
            "`backend/src/app.ts` — montage routes & middleware.",
            "`backend/prisma/schema.prisma` — modèle de données.",
            "`backend/package.json` — stack & scripts.",
            "`docker-compose*.yml` — topologie Docker.",
            "`.github/workflows/ci.yml` — pipeline.",
        ],
        "Front & doc",
        [
            "`package.json` (racine PWA deps).",
            "`frontend/package.json` (Expo).",
            "`frontend/src/api/client.ts` — stratégie API mobile.",
            "`docs/ARCHITECTURE.md`.",
            "`vite.config.js` — build PWA.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
