# -*- coding: utf-8 -*-
"""
AfriWonder — Génération complète de la soutenance (contenu utilisateur + illustrations).

Sortie : AfriWonder_Soutenance.pptx

Usage :
  cd docs/soutenance
  pip install python-pptx
  python generate_soutenance_complet.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(11, 17, 33)
BG_CARD = RGBColor(21, 30, 53)
EDGE = RGBColor(56, 189, 248)
TEXT = RGBColor(248, 250, 252)
TEXT_MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(56, 189, 248)
AMBER = RGBColor(245, 158, 11)
GREEN = RGBColor(52, 211, 153)
FONT_FALLBACK = "Calibri"

ROOT = Path(__file__).resolve().parent
REPO_ROOT = ROOT.parent.parent
OUT = ROOT / "AfriWonder_Soutenance.pptx"

# Captures : chemins relatifs à la racine du dépôt (premier existant utilisé pour chaque étiquette)
_SCREENSHOT_PATHS: dict[str, list[tuple[str, ...]]] = {
    "Page d’accueil / Landing": [
        ("public", "og-image.png"),
        ("frontend", "assets", "images", "app-image.png"),
    ],
    "Feed vidéo": [("public", "landing-feed-ui.png")],
    "Marketplace": [("public", "landing-market-ui.png")],
    "Wallet": [],
    "Messagerie (AfriChat)": [],
    "Panneau administrateur": [],
    "Application mobile (Expo)": [
        ("frontend", "assets", "images", "app-image.png"),
        ("public", "landing-live-ui.png"),
    ],
}
_SCREEN_GRID_ORDER: list[str] = [
    "Page d’accueil / Landing",
    "Feed vidéo",
    "Marketplace",
    "Wallet",
    "Messagerie (AfriChat)",
    "Panneau administrateur",
]


def _font(run, name: str = FONT_FALLBACK) -> None:
    run.font.name = name


def _bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _watermark(slide) -> None:
    el = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.9), Inches(1.05), Inches(5.5), Inches(5.3))
    el.fill.solid()
    el.fill.fore_color.rgb = RGBColor(14, 20, 38)
    el.line.fill.background()


def slide_blank(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _watermark(s)
    return s


def tit(s, txt: str, *, top=0.32, size=28, color=CYAN):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    _font(p.runs[0])


def notes(s, txt: str) -> None:
    s.notes_slide.notes_text_frame.text = txt.strip()


def box(
    s,
    left,
    top,
    w,
    h,
    text,
    *,
    accent=EDGE,
    fill=BG_CARD,
    bold=True,
    size=11,
    left_align=False,
    center=False,
):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = accent
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP if left_align else MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    p0.text = lines[0]
    p0.font.size = Pt(size)
    p0.font.bold = bold
    p0.font.color.rgb = TEXT
    _font(p0.runs[0])
    al = PP_ALIGN.LEFT if left_align else (PP_ALIGN.CENTER if center else PP_ALIGN.LEFT)
    p0.alignment = al
    for ln in lines[1:]:
        pr = tf.add_paragraph()
        pr.text = ln if ln.strip() else "\u00a0"
        pr.font.size = Pt(max(size - 1, 9))
        pr.font.bold = False
        pr.font.color.rgb = TEXT_MUTED
        if pr.runs:
            _font(pr.runs[0])
        pr.alignment = al
        pr.space_before = Pt(3)


def arrow_h(s, x1, y, x2):
    if x2 <= x1:
        return
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = EDGE
    c.line.width = Pt(1.2)


def _first_repo_image(label: str) -> Path | None:
    for parts in _SCREENSHOT_PATHS.get(label, []):
        cand = REPO_ROOT.joinpath(*parts)
        if cand.is_file():
            return cand
    return None


def _add_picture_fit(s, path: Path, left_in: float, top_in: float, max_w_in: float, max_h_in: float):
    pic = s.shapes.add_picture(str(path), Inches(left_in), Inches(top_in))
    mx = float(Inches(max_w_in))
    mh = float(Inches(max_h_in))
    w = float(pic.width)
    h = float(pic.height)
    if w <= 0 or h <= 0:
        return pic
    scale = min(mx / w, mh / h, 1.0)
    if scale < 1:
        pic.width = int(round(w * scale))
        pic.height = int(round(h * scale))
    return pic


def _screen_placeholder(s, left: float, top: float, w: float, h: float, label: str) -> None:
    ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    ph.adjustments[0] = 0.06
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(15, 23, 42)
    ph.line.color.rgb = CYAN
    ph.line.width = Pt(1.25)
    ph.line.dash_style = 1
    tf = ph.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"À insérer\ncapture\n\n{label}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN
    _font(p.runs[0])
    p.alignment = PP_ALIGN.CENTER


def slide_objectifs(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Objectifs du projet", size=28)
    box(
        s,
        0.68,
        1.22,
        6.12,
        5.25,
        "Objectifs généraux\n\n"
        "• Créer une super-app africaine multi-services\n"
        "• Centraliser les services numériques dans une seule plateforme\n"
        "• Optimiser l’accessibilité mobile et web",
        accent=CYAN,
        bold=False,
        size=14,
        left_align=True,
    )
    box(
        s,
        6.95,
        1.22,
        6.1,
        5.25,
        "Objectifs techniques\n\n"
        "• Concevoir une architecture scalable\n"
        "• Assurer la sécurité des transactions\n"
        "• Permettre le temps réel et la haute disponibilité",
        accent=AMBER,
        bold=False,
        size=14,
        left_align=True,
    )
    foot = s.shapes.add_textbox(Inches(0.72), Inches(6.72), Inches(12), Inches(0.62)).text_frame.paragraphs[0]
    foot.text = (
        "Cette slide relie la vision métier aux exigences d’ingénierie avant la problématique et les besoins détaillés."
    )
    foot.font.size = Pt(11)
    foot.font.color.rgb = TEXT_MUTED
    foot.alignment = PP_ALIGN.CENTER
    _font(foot.runs[0])
    notes(
        s,
        "Insistez : objectifs généraux = valeur utilisateur et territoire ; objectifs techniques = "
        "qualités non fonctionnelles (scalabilité, sécurité, disponibilité, temps réel).",
    )


def slide_methodologie(prs: Presentation) -> None:
    slide_bullets(
        prs,
        "Méthodologie de travail",
        [
            "Méthode Agile / Scrum (itérations courtes, livrables incrémentaux, revues régulières)",
            "Gestion des versions avec Git / GitHub (branches thématiques, Pull Requests)",
            "Découpage en sprints avec objectifs livrables (features + correctifs)",
            "Suivi des tâches (issues / backlog GitHub ou équivalent : À faire · En cours · Terminé)",
            "Workflow CI/CD : lint, typecheck, tests automatisés, garde qualité avant merge et déploiement",
        ],
        "À l’ENSAK MGSI : montrer le chaînon Organisation → Outils → Qualité logicielle.\n"
        "Évoquez en une phrase vos rituels Scrum (daily courte, sprint review, rétrospective si applicable).",
    )


def slide_app_screens_grid(prs: Presentation) -> None:
    """6 vignettes Web/PWA (+ placeholders). Insère automatiquement les PNG présents sous public/."""
    s = slide_blank(prs)
    tit(s, "Captures d’écran — partie 1 / 2 (Web & modules)", size=24)
    sub = (
        "Visuels issus du dépôt lorsqu’ils existent (dossier public/). Complétez Wallet, Messagerie et Admin "
        "avec vos captures réelles (Win+Shift+S ou outil téléphone)."
    )
    cap = s.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(12), Inches(0.82)).text_frame.paragraphs[0]
    cap.text = sub
    cap.font.size = Pt(11)
    cap.font.color.rgb = TEXT_MUTED
    cap.alignment = PP_ALIGN.LEFT
    _font(cap.runs[0])

    cols, rows = 3, 2
    margin_left, gap_x, gap_y = 0.55, 0.22, 0.52
    col_w = (13.333 - margin_left * 2 - gap_x * (cols - 1)) / cols
    cap_h = 0.28
    img_top_pad = 0.08
    row_img_h = 2.42
    y0 = 1.92

    idx = 0
    for r in range(rows):
        for c in range(cols):
            if idx >= len(_SCREEN_GRID_ORDER):
                break
            label = _SCREEN_GRID_ORDER[idx]
            left = margin_left + c * (col_w + gap_x)
            top = y0 + r * (cap_h + img_top_pad + row_img_h + gap_y)

            lbl = s.shapes.add_textbox(Inches(left), Inches(top), Inches(col_w), Inches(cap_h)).text_frame
            lp = lbl.paragraphs[0]
            lp.text = label
            lp.font.bold = True
            lp.font.size = Pt(11)
            lp.font.color.rgb = CYAN
            _font(lp.runs[0])

            pic_top = top + cap_h + img_top_pad
            img_path = _first_repo_image(label)
            if img_path:
                try:
                    _add_picture_fit(s, img_path, left, pic_top, col_w, row_img_h)
                except Exception:
                    _screen_placeholder(s, left, pic_top, col_w, row_img_h, label)
            else:
                _screen_placeholder(s, left, pic_top, col_w, row_img_h, label)
            idx += 1

    notes(
        s,
        "Dites au jury ce qui est une maquette/marketing PNG du repo vs une capture fonctionnelle live.\n"
        "Complétez les cases en pointillés avant la défense.",
    )


def slide_app_screens_mobile(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Captures d’écran — partie 2 / 2 (Application mobile Expo)", size=24)
    label = "Application mobile (Expo)"
    lbl = s.shapes.add_textbox(Inches(0.72), Inches(1.12), Inches(12), Inches(0.38)).text_frame.paragraphs[0]
    lbl.text = label + " — feed, création vidéo, notifications, etc."
    lbl.font.size = Pt(12)
    lbl.font.bold = True
    lbl.font.color.rgb = CYAN
    _font(lbl.runs[0])

    img_path = _first_repo_image(label)
    left_margin = 3.85
    max_w = 5.62
    max_h = 5.62
    top_pic = 1.68
    if img_path:
        try:
            _add_picture_fit(s, img_path, left_margin, top_pic, max_w, max_h)
        except Exception:
            _screen_placeholder(s, left_margin, top_pic, max_w, max_h, label)
    else:
        _screen_placeholder(s, left_margin, top_pic, max_w, max_h, label)

    legend = (
        "Conseil : insérez ici une capture du flux principal sur émulateur ou appareil (barre de statut visible, "
        "version lisible)."
    )
    foot = s.shapes.add_textbox(Inches(0.72), Inches(6.55), Inches(12), Inches(0.82)).text_frame.paragraphs[0]
    foot.text = legend
    foot.font.size = Pt(11)
    foot.font.color.rgb = TEXT_MUTED
    foot.alignment = PP_ALIGN.CENTER
    _font(foot.runs[0])
    notes(
        s,
        "Si vous projetez depuis un téléphone, cette slide peut être remplacée par une véritable démo live ; "
        "sinon conservez une capture nette avec identité AfriWonder visible.",
    )


def slide_bullets(prs: Presentation, title: str, bullets: list[str], note: str) -> None:
    s = slide_blank(prs)
    tit(s, title, size=26)
    tb = s.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.8), Inches(5.95))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(14 if len(bullets) < 8 else 12.5)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(6)
        _font(p.runs[0])
    notes(s, note)


def slide_cover(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Soutenance de Projet", top=1.05, size=34)
    sub = s.shapes.add_textbox(Inches(0.72), Inches(1.85), Inches(12), Inches(1.55)).text_frame
    sp = sub.paragraphs[0]
    sp.text = "AfriWonder Super-App"
    sp.font.size = Pt(26)
    sp.font.bold = True
    sp.font.color.rgb = TEXT
    _font(sp.runs[0])
    sp.alignment = PP_ALIGN.CENTER

    sp2 = sub.add_paragraph()
    sp2.text = "Écosystème numérique multi-services"
    sp2.font.size = Pt(16)
    sp2.font.color.rgb = CYAN
    sp2.space_before = Pt(10)
    _font(sp2.runs[0])
    sp2.alignment = PP_ALIGN.CENTER

    meta = (
        "Candidat : FANE ABDOULAYE · Élève Ingénieur\n"
        "Encadrant : Pr. HAMZA KHALFI\n"
        "Filière : Management et Gouvernance des Systèmes d’Information (MGSI)\n"
        "École Nationale des Sciences Appliquées de Khouribga (ENSAK)"
    )
    box(s, 1.95, 3.95, 9.42, 2.85, meta, accent=EDGE, bold=False, size=13.5, left_align=True)

    footer = s.shapes.add_textbox(Inches(0.72), Inches(6.82), Inches(12), Inches(0.45)).text_frame.paragraphs[0]
    footer.text = "AfriWonder — Projet d’ingénierie MGSI · Présentation complète avec illustrations"
    footer.font.size = Pt(11)
    footer.font.color.rgb = TEXT_MUTED
    footer.alignment = PP_ALIGN.CENTER
    _font(footer.runs[0])

    notes(
        s,
        "Accueillez le jury avec le titre exact du projet, votre identité ENSAK/MGSI, et l’encadrant.\n"
        "Annoncez la durée (ex. 12–15 minutes) puis le déroulé.",
    )


def slide_toc(prs: Presentation) -> None:
    items = [
        "Introduction et contexte du projet",
        "Objectifs du projet",
        "Méthodologie de travail",
        "Analyse fonctionnelle et métiers",
        "Architecture système et stack technologique",
        "Conception de la base de données",
        "Développement mobile et PWA",
        "Sécurité, performance et DevOps",
        "Conclusion et perspectives",
    ]
    slide_bullets(
        prs,
        "Sommaire de la présentation",
        items,
        "Structurez votre oral : du contexte aux objectifs et à la méthode, puis fonctionnel, technique, "
        "preuves (captures / CI) et conclusion.",
    )
    last = prs.slides[-1]
    y0 = 1.82
    row = 0.44
    for i in range(len(items)):
        box(last, 11.72, y0 + i * row, 0.5, 0.44, str(i + 1), accent=AMBER, size=10, center=True)


def illustration_arch_layers(prs: Presentation, title_note: tuple[str, str]) -> None:
    title, note = title_note
    s = slide_blank(prs)
    tit(s, title, size=26)
    # Couches empilées
    layers = [
        ("Clients", "PWA (Vite/React) · Mobile (Expo/RN)", 1.35, CYAN),
        ("Communication", "REST JWT · WebSockets Socket.IO", 2.52, EDGE),
        ("Backend", "Express · Routes / Services · Zod", 3.72, AMBER),
        ("Persistance", "Prisma · PostgreSQL", 4.9, GREEN),
        ("Cache temps réel", "Redis · rate limit · adapter sockets", 6.08, CYAN),
    ]
    lx, w = 2.95, 7.42
    for lab, sub, top, ac in layers:
        box(s, lx, top, w, 0.88, f"{lab}\n{sub}", accent=ac, bold=True, size=11.5, left_align=True)
        if top < 6.1:
            ln = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(lx + w / 2),
                Inches(top + 0.89),
                Inches(lx + w / 2),
                Inches(top + 1.0),
            )
            ln.line.color.rgb = EDGE
            ln.line.width = Pt(1)
    illustration_hint = (
        "Lecture : flux descendant (clients → API → données). "
        "Socket.IO permet le temps réel ; Redis aide à scaler horizontalement quand plusieurs instances backend tournent."
    )
    notes(s, f"{note}\n\n{illustration_hint}")


def illustration_user_hub(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Schéma : User comme hub métier", size=26)
    hx, hy = 5.92, 2.92
    hub = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(hx), Inches(hy), Inches(2.15), Inches(1.05))
    hub.fill.solid()
    hub.fill.fore_color.rgb = BG_CARD
    hub.line.color.rgb = CYAN
    hub.line.width = Pt(1.5)
    hf = hub.text_frame
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = hf.paragraphs[0]
    hp.text = "USER"
    hp.font.bold = True
    hp.font.size = Pt(17)
    hp.font.color.rgb = TEXT
    _font(hp.runs[0])
    hp.alignment = PP_ALIGN.CENTER

    sats = [
        ("Vidéos", 0.88, 1.35),
        ("Order", 9.75, 1.35),
        ("Wallet", 0.82, 3.92),
        ("Stories", 9.92, 3.92),
        ("Messagerie", 0.86, 5.82),
        ("Marketplace", 9.72, 5.82),
    ]
    cx = hx + 1.075
    cy = hy + 0.525
    for lab, lx, ly in sats:
        box(s, lx, ly, 2.42, 0.62, lab, accent=AMBER, bold=True, size=12)
        ln = s.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(lx + 1.21),
            Inches(ly + 0.31),
            Inches(cx),
            Inches(cy),
        )
        ln.line.color.rgb = RGBColor(71, 85, 105)
        ln.line.width = Pt(1)

    cap = s.shapes.add_textbox(Inches(0.72), Inches(6.55), Inches(12), Inches(0.72)).text_frame.paragraphs[0]
    cap.text = "Illustration relationnelle : un utilisateur peut cumuler rôles (créateur, acheteur, prestataire) selon modules."
    cap.font.size = Pt(11.5)
    cap.font.color.rgb = TEXT_MUTED
    cap.alignment = PP_ALIGN.CENTER
    _font(cap.runs[0])

    notes(
        s,
        "Expliquez au jury que PostgreSQL impose l’intégrité référentielle : tout module sensible (paiement, livraison) "
        "reste relié à une identité stable (User).\nLes traits sont une simplification du schéma Prisma réel (~100 modèles).",
    )


def illustration_realtime(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Illustration : temps réel (Socket.IO)", size=26)
    box(s, 0.85, 2.05, 2.35, 1.85, "Client A\nMessagerie", accent=CYAN)
    arrow_h(s, 3.22, 2.95, 4.96)
    box(s, 4.98, 2.42, 3.92, 1.12, "Serveur Socket.IO\nrooms · auth JWT", accent=EDGE)
    arrow_h(s, 8.92, 2.95, 10.72)
    box(s, 10.74, 2.05, 2.05, 1.85, "Client B\nPrésence", accent=AMBER)
    box(s, 5.42, 4.42, 3.52, 0.78, "Redis adapter\n(multi-instance)", accent=GREEN, bold=False)

    ln = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.0),
        Inches(3.54),
        Inches(7.0),
        Inches(4.42),
    )
    ln.line.color.rgb = TEXT_MUTED
    ln.line.dash_style = 2

    bullets = (
        "Événements typiques : message, typing, notification, géolocalisation ride/livraison (selon modules activés)."
    )
    footer = s.shapes.add_textbox(Inches(0.72), Inches(5.55), Inches(12), Inches(1.05)).text_frame.paragraphs[0]
    footer.text = bullets
    footer.font.size = Pt(12)
    footer.font.color.rgb = TEXT_MUTED
    footer.alignment = PP_ALIGN.CENTER
    _font(footer.runs[0])

    notes(
        s,
        "Précisez : la messagerie et les salons utilisent Socket.IO.\nEn production avec plusieurs backends, REDIS_URL active "
        "l’adaptateur Redis pour diffuser les événements entre instances.",
    )


def illustration_payments(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Illustration : flux paiement (vue haut niveau)", size=25)
    steps = [
        ("1 — Initiation", "POST sécurisé\nidempotence · risque/KYC selon périmètre", 1.42),
        ("2 — Provider", "Stripe / OM / autres\nvérif signatures", 3.92),
        ("3 — Webhook", "body brut · secret\nreplay protection", 6.52),
        ("4 — Confirmation", "mise à jour commande /\nwallet / abonnements", 8.92),
    ]
    x = 1.42
    for lab, desc, lx in [(t[0], t[1], float(t[2])) for t in steps]:
        box(s, lx, 2.2, 2.38, 1.92, f"{lab}\n{desc}", accent=EDGE, bold=False, size=11, left_align=True)
        if lx < 8.4:
            arrow_h(s, lx + 2.41, 3.05, lx + 2.73)

    footer = s.shapes.add_textbox(Inches(0.72), Inches(5.25), Inches(12), Inches(1.45)).text_frame.paragraphs[0]
    footer.text = (
        "Le projet encode plusieurs passerelles dans `payments.routes` : "
        "le jury peut être informé du positionnement régional ET international."
    )
    footer.font.size = Pt(11.5)
    footer.font.color.rgb = TEXT_MUTED
    footer.alignment = PP_ALIGN.CENTER
    _font(footer.runs[0])

    notes(
        s,
        "Insistez sur la conformité sécuritaire : webhook en body brut, vérifications, et désactivation du « trust webhook » "
        "en production.\n Mentionnez brièvement Stripe + paiements mobiles locaux disponibles dans l’implémentation.",
    )


def illustration_devops(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Illustration : chaîne DevOps simplifiée", size=26)
    chain = [
        ("Git", 0.85),
        ("CI\nGitHub Actions", 2.85),
        ("Build / tests", 5.12),
        ("Docker Compose", 7.45),
        ("Nginx TLS", 9.85),
    ]
    for i, (lab, left) in enumerate(chain):
        box(s, left, 2.25, 1.85, 1.45, lab, accent=CYAN if i == 0 else EDGE, size=11.5)
        if i < len(chain) - 1:
            nxt = chain[i + 1][1]
            arrow_h(s, left + 1.88, 2.95, nxt - 0.2)

    box(s, 1.28, 4.15, 2.55, 0.78, "3× backend\n(haute dispo — compose)", accent=AMBER)
    box(s, 4.25, 4.15, 2.25, 0.78, "Postgres 15", accent=GREEN)
    box(s, 6.85, 4.15, 2.18, 0.78, "Redis 7", accent=CYAN)
    box(s, 9.08, 4.15, 3.08, 0.78, "Certbot TLS", accent=EDGE)

    notes(
        s,
        "À l’oral : la réplication compose est un objectif documenté dans le repo ; précisez la plateforme de déploiement réelle "
        "(serveur VPS, Render, etc.) pour rester vérifiable.",
    )


def slide_thanks(prs: Presentation) -> None:
    s = slide_blank(prs)
    tit(s, "Merci de votre attention", top=2.45, size=36, color=TEXT)
    tb = s.shapes.add_textbox(Inches(0.92), Inches(3.35), Inches(11.58), Inches(1.5)).text_frame
    q = tb.paragraphs[0]
    q.text = "Questions ?"
    q.font.size = Pt(26)
    q.font.bold = False
    q.font.color.rgb = CYAN
    _font(q.runs[0])
    q.alignment = PP_ALIGN.CENTER

    sub = tb.add_paragraph()
    sub.text = (
        "AfriWonder · Super-App régionale · PWA · Mobile Expo · Backend Express · PostgreSQL · "
        "Sécurité & observabilité"
    )
    sub.font.size = Pt(13)
    sub.font.color.rgb = TEXT_MUTED
    sub.space_before = Pt(18)
    _font(sub.runs[0])
    sub.alignment = PP_ALIGN.CENTER

    notes(s, "Clôture courte : résultat, limites assumées, ouverture industrialisation.")


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)

    slide_toc(prs)

    slide_bullets(
        prs,
        "Introduction au projet AfriWonder",
        [
            "Genèse : répondre à la fragmentation numérique en Afrique",
            "Vision : une « Super-App » unifiée (Social, Commerce, Finance, Services)",
            "Ambition : centraliser l’expérience dans un écosystème régional cohérent",
            "Valeur ajoutée : accessibilité, monétisation locale, services de proximité",
        ],
        "Donnez le « pourquoi » avant le « comment » : un seul compte, des parcours bout-en-bout, "
        "une vision produit alignée sur les usages mobile money et la faible bande passante.",
    )

    slide_objectifs(prs)
    slide_methodologie(prs)

    slide_bullets(
        prs,
        "Problématique et enjeux stratégiques",
        [
            "Fragmentation des services mobiles en Afrique",
            "Difficultés d’accès aux infrastructures de paiement internationales",
            "Contraintes de connectivité et de performance des terminaux",
            "Nécessité d’une gouvernance centralisée des données (MGSI)",
        ],
        "Reliez chaque enjeu à un choix d’architecture : API unifiée, cache, compression, "
        "passerelles de paiement locales, traçabilité et contrôle des accès.",
    )

    slide_bullets(
        prs,
        "Analyse des besoins fonctionnels",
        [
            "Modules principaux : feed vidéo, marketplace, messagerie temps réel",
            "Services intégrés : wallet, transport, santé, tontines",
            "Profils : créateurs, vendeurs, prestataires, administrateurs",
            "Transverse : traduction, chatbots, notifications push",
        ],
        "Présentez un découpage par domaines : social, commerce, finance, services locaux, "
        "puis les briques transverses (i18n, assistance, alertes).",
    )

    slide_bullets(
        prs,
        "Cartographie des cas d’utilisation",
        [
            "Social : publication et interaction sur le contenu vidéo",
            "Marketplace : cycle annonce → panier → paiement sécurisé",
            "Wallet : solde, historique, règles de sécurité (PIN / limites selon implémentation)",
            "Gouvernance : modération et administration globale",
        ],
        "Insistez sur la complétude du parcours : ce n’est pas une vitrine, c’est un système transactionnel.",
    )

    illustration_arch_layers(
        prs,
        (
            "Architecture logique du système",
            "Modèle client-serveur distribué, API REST sécurisée, WebSockets ; "
            "découpage modulaire par domaine ; couches routes / services / données.",
        ),
    )

    slide_bullets(
        prs,
        "Stack technologique : choix d’efficacité",
        [
            "Typescript côté backend et sur une grande partie des clients (alignement d’équipe)",
            "Backend : Node.js + Express",
            "Web PWA : React 18 + Vite",
            "Mobile : React Native + Expo (SDK 54+)",
        ],
        "Nuance orale : le dépôt peut contenir aussi du JS historique ; l’objectif est la cohérence TypeScript sur le code critique.",
    )

    slide_bullets(
        prs,
        "Architecture backend et API",
        [
            "Routage structuré par domaine métier (social, marketplace, wallet, services…)",
            "Validation Zod sur les entrées sensibles",
            "Middlewares : JWT, rate limiting, Helmet, CORS",
            "Médias : Multer + traitements (ex. Sharp / transcodage selon chemins)",
        ],
        "Expliquez la séparation des responsabilités : route fine, service métier, accès Prisma, "
        "journalisation et gestion d’erreurs centralisée.",
    )

    slide_bullets(
        prs,
        "Modélisation des données avec Prisma",
        [
            "ORM Prisma pour productivité et typage fort",
            "PostgreSQL relationnel, schéma étendu (≈ 100 modèles dans le schéma Prisma du dépôt)",
            "Migrations automatisées",
            "Optimisation des requêtes pour volumes importants (index, requêtes ciblées)",
        ],
        "Le schéma massif est un atout (richesse fonctionnelle) et un risque (onboarding) : "
        "dites comment vous le documentez (Swagger, conventions de nommage).",
    )

    illustration_user_hub(prs)

    slide_bullets(
        prs,
        "Structure du schéma (focus métier)",
        [
            "User : pivot vidéos, wallet, commandes, rôles",
            "Marketplace : produits, catégories, panier, stocks",
            "Wallet : transactions, soldes, historique",
            "Social : abonnements, commentaires, likes, stories",
        ],
        "Reprenez le schéma illustré : cohérence transactionnelle et contraintes d’intégrité.",
    )

    slide_bullets(
        prs,
        "Stratégie frontend : PWA vs mobile natif (Expo)",
        [
            "PWA : reach web, SEO, déploiement rapide, expérience légère",
            "Mobile : caméra, notifications push, modules natifs, usage terrain",
            "Partage des principes d’intégration API (même contrat backend)",
            "UI web : Tailwind + composants ; mobile : design system RN cohérent avec la marque",
        ],
        "Ne promettez pas « un seul design system pixel-perfect » : dites « cohérence UX et patterns partagés ».",
    )

    slide_bullets(
        prs,
        "État serveur et expérience utilisateur",
        [
            "TanStack Query pour cache, re-fetch et résilience réseau",
            "Zustand pour état global léger (session, préférences)",
            "Persistance locale (stockage sécurisé côté mobile selon modules)",
            "Optimistic UI sur interactions sociales (likes, commentaires) lorsque pertinent",
        ],
        "Soulignez le bénéfice en Afrique : moins de requêtes inutiles, meilleure UX sur connexion instable.",
    )

    illustration_realtime(prs)

    slide_bullets(
        prs,
        "Communication temps réel (détail)",
        [
            "Socket.IO pour messagerie instantanée",
            "Salons / rooms et présence en ligne",
            "Notifications événementielles (social et commerce)",
            "Scalabilité : adaptateur Redis quand plusieurs instances backend",
        ],
        "Référez-vous au schéma : authent JWT sur handshake, rooms par utilisateur/conversation.",
    )

    illustration_payments(prs)

    slide_bullets(
        prs,
        "Services de paiement",
        [
            "Passerelles : Stripe + solutions locales (ex. Orange Money, Wave, Moov…) selon périmètre déployé",
            "Webhooks avec vérifications (signatures / règles anti-abus)",
            "Workflow typique : init → callback/provider → mise à jour commande ou wallet",
            "Traçabilité financière : logs d’audit, idempotency sur init sensible",
        ],
        "Soyez factuel sur ce qui est activé sur votre environnement de démo (clés prod vs sandbox).",
    )

    slide_bullets(
        prs,
        "Traduction & intelligence artificielle",
        [
            "Traduction LibreTranslate (+ repli MyMemory selon routes)",
            "Chatbots pour assistance contextualisée (données métier)",
            "STT / sous-titres : chantier en cours (Whisper/OpenAI ou pipeline interne)",
            "Modération assistée par règles + listes de contrôle",
        ],
        "Dites explicitement où c’est encore partiel pour éviter les questions surprises du jury.",
    )

    slide_bullets(
        prs,
        "Sécurité du système d’information",
        [
            "JWT access + refresh, rotation/blacklist selon implémentation",
            "XSS / CSRF : protections middleware + bonnes pratiques cookies",
            "Injection SQL : Prisma (requêtes paramétrées)",
            "Proxy média anti-SSRF : liste blanche de domaines",
            "Secrets : variables d’environnement, pas de credentials en repo",
        ],
        "Ajoutez un exemple oral : tentative d’abus webhook rejetée faute de signature valide.",
    )

    slide_bullets(
        prs,
        "Performance et optimisations",
        [
            "Code splitting et lazy loading côté frontend web",
            "Compression HTTP + stratégie de cache Redis côté API",
            "Vidéos : streaming par fragments lorsque disponible ; adaptation réseaux lents",
            "Observabilité : Sentry + exposition métriques /metrics format Prometheus",
        ],
        "Connectez perf et contexte Mali/Afrique : latence, coût data, cold start.",
    )

    illustration_devops(prs)

    slide_bullets(
        prs,
        "Stratégie DevOps et déploiement",
        [
            "Docker + Docker Compose pour reproductibilité",
            "Haute disponibilité : plusieurs réplicas backend (fichiers compose)",
            "Nginx reverse proxy + TLS (Let’s Encrypt / Certbot)",
            "PostgreSQL 15 et Redis 7",
        ],
        "Précisez si votre démo est locale, staging cloud, ou production — même architecture, paramètres différents.",
    )

    slide_bullets(
        prs,
        "Automatisation et qualité (CI/CD)",
        [
            "GitHub Actions : intégration continue",
            "Contrôles automatiques lint + typecheck",
            "Tests backend (Jest) et frontend (Vitest) + sélection E2E Playwright sur le repo",
            "Contraintes de taille de PR pour garder une review efficace",
        ],
        "Le jury MGSI adore la boucle Qualité→Livraison : montrez où ça se voit dans le dépôt GitHub.",
    )

    slide_bullets(
        prs,
        "Périmètre livré vs feuille de route",
        [
            "Livré : socle Super-App, API riche, clients web/mobile, temps réel, observabilité, CI",
            "Partiel / variable selon configuration : IA vocale/STT bout-en-bout production",
            "Roadmap courte : modularisation, charge, feature flags pilotés produit",
        ],
        "Cette slide « protège » votre crédibilité : vous projetez mais vous ne cachez pas l’état réel.",
    )

    slide_bullets(
        prs,
        "Preuves techniques (vérifiables dans le dépôt)",
        [
            "Schéma Prisma : environ 100 modèles PostgreSQL",
            "Socket.IO + adaptateur Redis optionnel pour multi-instance",
            "Endpoints de santé/métriques et intégration Sentry backend/mobile",
            "Nombre important de routes domaine + tests automatisés en CI",
        ],
        "Utilisez « vérifiable » : le jury peut demander où — vous pointez dossiers/backend/CI.",
    )

    slide_app_screens_grid(prs)
    slide_app_screens_mobile(prs)

    slide_bullets(
        prs,
        "Défis techniques & solutions",
        [
            "Complexité du schéma Prisma massif → conventions, migrations, modularisation progressive",
            "Lecture vidéo sur réseaux instables → politiques buffer, qualités adaptatives",
            "Cohérence auth PWA/mobile → même contrat JWT, refresh, stockage mobile sécurisé",
            "Proxy média SSRF potentiel → liste blanche stricte & revues sécurité",
        ],
        "Pour chaque défi : symptôme, cause, contre-mesure, résultat partiel/complet.",
    )

    slide_bullets(
        prs,
        "Limites actuelles assumées",
        [
            "Monolithe encore dominant (scalabilité verticalisée puis découpage progressif)",
            "Dépendance à fournisseurs IA/traduction pour certains SLA",
            "STT/Speech-to-text : placeholders ou finalisation environnement-specific",
            "Schéma DB dense : courbe d’apprentissage pour nouvelles équipes",
        ],
        "Terminez cette slide par une phrase « ce n’est pas un échec, c’est une dette projetée avec plan ».",
    )

    slide_bullets(
        prs,
        "Améliorations futures et évolutions",
        [
            "Microservices ou services autonomes pour domaines critiques (paiements, médias)",
            "IA/STT/NLP finalisées et autosuffisantes quand données & consentement permettent",
            "Tests de montée en charge industriels au-delà des scripts projet",
            "Feature flags pour déploiement progressif sans casser les clients",
        ],
        "Clarifiez que la transition architecture est coûteuse : elle se fait après traction et métriques.",
    )

    slide_bullets(
        prs,
        "Compétences acquises & bilan personnel",
        [
            "Cycle de vie full-stack réel : conception → implémentation → tests → exploitation",
            "Gouvernance SI : données, sécurité, conformité fonctionnelle (angle MGSI)",
            "Architectures évolutives : résilience réseaux pauvres, supervision",
            "Culture qualité industrielle : PR, automatisation CI, documentation",
        ],
        "Reprendre 2 compétences en lien avec le référentiel ENSAK MGSI.",
    )

    slide_bullets(
        prs,
        "Conclusion générale",
        [
            "Réponse techno concrète au marché africain (fragmentation vs plateforme unifiée)",
            "Projet complet : innovation UX + rigueur ingénierie",
            "Validation académique des compétences MGSI",
            "Ouverture : industrialisation progressive et mesure terrain",
        ],
        "Une phrase forte : AfriWonder n’est pas un prototype isolé mais une base industrielle évolutive.",
    )

    slide_thanks(prs)
    return prs


def main() -> None:
    deck = build()
    deck.save(OUT)
    print(f"OK : présentation générée — {OUT}")


if __name__ == "__main__":
    main()
