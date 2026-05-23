# -*- coding: utf-8 -*-
"""
AfriWonder — soutenance PowerPoint (projection publique vs préparation perso).

Une vraie présentation : peu de texte à l’écran, visuels lisibles — pas de mémoires entiers ni fiches « pièges » projetés.

Sorties :
  • AfriWonder_Soutenance.pptx — À montrer devant la classe (visuels + Merci).
  • AfriWonder_Soutenance_PREP_JURY.pptx — Pour toi : mêmes visuels puis diapos ÉCRAN NEUTRE ;
    tout le markdown des mémoires est uniquement dans les NOTES présentateur.

Option : `--public-only` pour ne régénérer que le fichier public.

Thème : fond #0B1121 · accents cyan / ambre · #F8FAFC.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Thème utilisateur (#0B1121, #F8FAFC, #38BDF8 ≈ cyan-400, #F59E0B amber-500)
BG = RGBColor(11, 17, 33)
BG_CARD = RGBColor(21, 30, 53)
EDGE = RGBColor(56, 189, 248)
TEXT = RGBColor(248, 250, 252)
TEXT_MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(56, 189, 248)
AMBER = RGBColor(245, 158, 11)
GREEN = RGBColor(52, 211, 153)

FONT_UI = "Inter"
FONT_FALLBACK = "Calibri"

ROOT = Path(__file__).resolve().parent
OUT_PUBLIC = ROOT / "AfriWonder_Soutenance.pptx"
OUT_PREP = ROOT / "AfriWonder_Soutenance_PREP_JURY.pptx"

MEMOIRE_FILES: list[tuple[str, Path]] = [
    ("Mémoire audit soutenance (complet)", ROOT / "MEMOIRE_AUDIT_SOUTENANCE_COMPLET.md"),
    ("Mémoire ingénieur (complet)", ROOT / "MEMOIRE_SOUTENANCE_INGENIEUR_COMPLET.md"),
    ("Annexe questions jury (techniques)", ROOT / "ANNEXE_QUESTIONS_JURY_TECHNIQUES.md"),
]

_HEADING = re.compile(r"^(#{1,3})\s+(.+)$")
_MAX_NOTE_CHARS = 26_000
_MAX_TITLE_LEN = 180

# Marque dans les NOTES si la section est à traiter comme sensible (pas de titre projeté dans le fichier PREP).
_PRIVATE_TITLE_HINTS_RE = re.compile(
    r"(piège|piege|priv[eé]|confidentiel|secret\s*jury|ne\s+pas\s+projeter"
    r"|appendice\s*a|question[s]?\s*piège|SECTION\s+PRIV|banque\s+élargie)",
    re.IGNORECASE,
)


def _run_font(run) -> None:
    try:
        run.font.name = FONT_UI
    except Exception:
        run.font.name = FONT_FALLBACK


def _notes(s, intention: str) -> None:
    s.notes_slide.notes_text_frame.text = intention.strip()


def _bg(s) -> None:
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def slide_new(prs: Presentation):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _bg(s)
    _watermark_map_hint(s)
    return s


def _watermark_map_hint(s) -> None:
    """Silhouette discrète (approx.) : cercles très sombres + note option image réelle."""
    el = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(3.9), Inches(1.1), Inches(5.5), Inches(5.2)
    )
    el.fill.solid()
    el.fill.fore_color.rgb = RGBColor(14, 20, 38)
    el.line.fill.background()


def _box(
    s,
    left: float,
    top: float,
    w: float,
    h: float,
    text: str,
    *,
    accent: RGBColor = EDGE,
    fill: RGBColor = BG_CARD,
    bold=True,
    title_size: int = 12,
    center=True,
    align_left: bool = False,
) -> None:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = accent
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if not align_left else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    parts = text.split("\n")
    p = tf.paragraphs[0]
    p.text = parts[0]
    p.font.size = Pt(title_size)
    p.font.bold = bold
    p.font.color.rgb = TEXT
    _run_font(p.runs[0])
    use_center = center and not align_left
    p.alignment = PP_ALIGN.CENTER if use_center else PP_ALIGN.LEFT
    for extra in parts[1:]:
        pr = tf.add_paragraph()
        pr.text = extra
        pr.font.size = Pt(max(title_size - 2, 9))
        pr.font.bold = False
        pr.font.color.rgb = TEXT_MUTED
        _run_font(pr.runs[0])
        pr.alignment = PP_ALIGN.LEFT if align_left else (PP_ALIGN.CENTER if use_center else PP_ALIGN.LEFT)
        pr.space_before = Pt(4)


def _title(s, text: str, *, top: float = 0.32, size: int = 30, color: RGBColor = CYAN) -> None:
    tb = s.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.0), Inches(0.85))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    _run_font(p.runs[0])


def _arrow_h(s, x1: float, y: float, x2: float) -> None:
    if x2 <= x1:
        return
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = EDGE
    c.line.width = Pt(1.25)


def _clean_md_title(raw: str) -> str:
    t = re.sub(r"<[^>]+>", "", raw).strip()
    if len(t) > _MAX_TITLE_LEN:
        t = t[: _MAX_TITLE_LEN - 1] + "…"
    return t


def parse_markdown_sections(md_text: str) -> list[tuple[str, str]]:
    """Découpe un Markdown en (titre de section, corps) selon # / ## / ###."""
    lines = md_text.splitlines()
    out: list[tuple[str, str]] = []
    buf: list[str] = []
    title: str | None = None
    in_code = False

    def flush() -> None:
        nonlocal title, buf
        if title is None:
            buf = []
            return
        body = "\n".join(buf).strip()
        buf = []
        ct = _clean_md_title(title)
        if not ct:
            return
        if not body:
            body = "(Section sans paragraphe sous-jacent dans le fichier Markdown — voir structure source.)"
        out.append((ct, body))

    for line in lines:
        st = line.strip()
        if st.startswith("```"):
            in_code = not in_code
            buf.append(line)
            continue
        if not in_code:
            hm = _HEADING.match(line)
            if hm:
                flush()
                title = hm.group(2).strip()
                continue
        buf.append(line)
    flush()
    return out


def _chunk_notes(text: str) -> list[str]:
    text = text.strip()
    if not text:
        return [""]
    chunks: list[str] = []
    s = text
    while len(s) > _MAX_NOTE_CHARS:
        cut = s.rfind("\n", 0, _MAX_NOTE_CHARS)
        if cut < _MAX_NOTE_CHARS // 4:
            cut = _MAX_NOTE_CHARS
        chunks.append(s[:cut].strip())
        s = s[cut:].strip()
    chunks.append(s)
    return chunks


def _bump_prep_face(face_n: list[int]) -> int:
    face_n[0] += 1
    return face_n[0]


def _slide_neutral_prep_face(prs: Presentation, face_n: list[int]):
    """Retourne (slide, numéro affiché)."""
    n = _bump_prep_face(face_n)
    s = slide_new(prs)
    _title(s, "Préparation soutenance", top=2.45, size=26)
    sub = s.shapes.add_textbox(Inches(0.85), Inches(3.15), Inches(11.6), Inches(1.25)).text_frame
    sp = sub.paragraphs[0]
    sp.text = f"Fiche {n}\n\nContenu réservé aux NOTES présentateur.\nNe pas projeter en groupe."
    sp.font.size = Pt(14)
    sp.font.color.rgb = TEXT_MUTED
    sp.alignment = PP_ALIGN.CENTER
    _run_font(sp.runs[0])
    return s, n


def _slide_prep_part_intro(prs: Presentation, face_n: list[int]) -> None:
    n = _bump_prep_face(face_n)
    s = slide_new(prs)
    _title(s, "Entraînement & mémoires", top=2.05, size=27)
    sub = s.shapes.add_textbox(Inches(0.92), Inches(2.95), Inches(11.45), Inches(2.3)).text_frame
    msg = (
        "À partir d’ici, le contenu détaillé des fichiers Markdown\n"
        "n’apparaît pas sur cette diapo — uniquement sous NOTES.\n\n"
        f"Balise personnelle · fiche {n}\n(ne pas diffuser en cours / promo)"
    )
    for i, line in enumerate(msg.split("\n")):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = CYAN if "Balise personnelle" in line else TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            _run_font(p.runs[0])
    _notes(
        s,
        "À utiliser avec Affichage > Notes ou mode Présentateur.\n"
        "Tu peux passer ces diapos en répétition sans exposer le texte sur le mur.",
    )


def append_markdown_prep_notes_only(prs: Presentation, face_n: list[int]) -> tuple[int, int]:
    """Diapos mémoires : écran neutre, texte intégral en NOTES. Retourne (sections_md, slides_créées)."""
    sections_seen = 0
    slides_made = 0
    for label, path in MEMOIRE_FILES:
        if not path.is_file():
            s = slide_new(prs)
            _title(s, "⚠ Fichier manquant", top=3.2, size=22)
            st = s.shapes.add_textbox(Inches(1.5), Inches(3.95), Inches(10), Inches(1.25)).text_frame.paragraphs[0]
            st.text = f"{path.name}\n(placez le .md dans docs/soutenance/)"
            st.font.size = Pt(13)
            st.font.color.rgb = AMBER
            _notes(s, f"Manquant : {path.resolve()}")
            slides_made += 1
            continue

        s_doc, doc_n = _slide_neutral_prep_face(prs, face_n)
        _notes(
            s_doc,
            f"Bloc document : {label}\nFichier : {path.resolve()}\n\n(Ouvre les diapos suivantes avec le panneau Notes.)",
        )
        slides_made += 1

        md_text = path.read_text(encoding="utf-8")
        for sec_title, sec_body in parse_markdown_sections(md_text):
            sections_seen += 1
            sensitive = _PRIVATE_TITLE_HINTS_RE.search(sec_title) is not None
            parts = _chunk_notes(sec_body)
            for i, chunk in enumerate(parts):
                s, n = _slide_neutral_prep_face(prs, face_n)
                header = [f"Fiche écran n°{n}", f"Source : {path.name}", ""]
                if sensitive:
                    header.append("[SECTION SENSIBLE — ne pas partager l’écran]")
                    header.append("")
                header.append(f"## {sec_title}")
                if len(parts) > 1:
                    header.append(f"(Notes découpées : partie {i + 1}/{len(parts)})")
                header.append("")
                header.append(chunk)
                s.notes_slide.notes_text_frame.text = "\n".join(header)
                slides_made += 1
    return sections_seen, slides_made


def _add_visual_deck_core(prs: Presentation) -> None:
    _slide_title_identity(prs)
    _slide_problems(prs)
    _slide_answer(prs)
    _slide_journey(prs)
    _slide_arch_global(prs)
    _slide_stack_table(prs)
    _slide_data_star(prs)
    _slide_security(prs)
    _slide_realtime(prs)
    _slide_devops(prs)
    _slide_ia(prs)
    _slide_difficulties(prs)


def _slide_title_identity(prs):
    s = slide_new(prs)
    logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(0.55), Inches(1.45), Inches(1.15))
    logo.fill.solid()
    logo.fill.fore_color.rgb = BG_CARD
    logo.line.color.rgb = CYAN
    logo.line.width = Pt(1)
    lf = logo.text_frame
    lf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lp = lf.paragraphs[0]
    lp.text = "A"
    lp.font.size = Pt(44)
    lp.font.bold = True
    lp.font.color.rgb = CYAN
    _run_font(lp.runs[0])
    lp.alignment = PP_ALIGN.CENTER

    hero = s.shapes.add_textbox(Inches(0.72), Inches(1.85), Inches(12.0), Inches(2.9)).text_frame
    hp = hero.paragraphs[0]
    hp.text = "AfriWonder"
    hp.font.size = Pt(40)
    hp.font.bold = True
    hp.font.color.rgb = TEXT
    _run_font(hp.runs[0])
    hp.alignment = PP_ALIGN.CENTER
    h2 = hero.add_paragraph()
    h2.text = "Super-app pour l’Afrique"
    h2.font.size = Pt(20)
    h2.font.color.rgb = CYAN
    _run_font(h2.runs[0])
    h2.alignment = PP_ALIGN.CENTER
    h2.space_before = Pt(8)
    h3 = hero.add_paragraph()
    h3.text = "PWA · Mobile · API"
    h3.font.size = Pt(17)
    h3.font.color.rgb = TEXT
    _run_font(h3.runs[0])
    h3.alignment = PP_ALIGN.CENTER
    h3.space_before = Pt(10)

    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.2), Inches(5.08), Inches(9.12), Inches(5.08))
    ln.line.color.rgb = EDGE

    meta = s.shapes.add_textbox(Inches(0.72), Inches(5.25), Inches(12), Inches(1.85)).text_frame
    mp0 = meta.paragraphs[0]
    mp0.text = "ABDOULAYE FANEL"
    mp0.font.size = Pt(22)
    mp0.font.bold = True
    mp0.font.color.rgb = TEXT
    _run_font(mp0.runs[0])
    mp0.alignment = PP_ALIGN.CENTER
    mp1 = meta.add_paragraph()
    mp1.text = "Encadré par Pr. HAMZA KALFI"
    mp1.font.size = Pt(16)
    mp1.font.color.rgb = TEXT_MUTED
    _run_font(mp1.runs[0])
    mp1.alignment = PP_ALIGN.CENTER
    mp1.space_before = Pt(6)
    mp2 = meta.add_paragraph()
    mp2.text = "ENSA Khouribga – MGSI"
    mp2.font.size = Pt(15)
    mp2.font.color.rgb = TEXT_MUTED
    _run_font(mp2.runs[0])
    mp2.alignment = PP_ALIGN.CENTER

    _notes(s, "Vous accueillez le jury et présentez l’intitulé du projet, votre nom et votre encadrant.")


def _slide_problems(prs):
    s = slide_new(prs)
    _title(s, "Problème et opportunité", top=0.38, size=26)
    items = [
        ("🔗 Connectivité instable", "Signal faible, latences, coûts data."),
        ("📱 Services éparpillés", "Expériences fragmentées sans fil conducteur."),
        ("💰 Monétisation difficile", "Peu de passerelles locales unifiées pour créateurs."),
    ]
    x0, w, gap, top = 0.92, 3.72, 0.52, 2.05
    for i, (t, sub) in enumerate(items):
        left = x0 + i * (w + gap)
        _box(s, left, top, w, 1.55, t, accent=CYAN if i % 2 == 0 else AMBER, title_size=14)
        st = s.shapes.add_textbox(Inches(left + 0.15), Inches(top + 1.55), Inches(w - 0.3), Inches(1.95)).text_frame
        sp = st.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(11.5)
        sp.font.color.rgb = TEXT_MUTED
        sp.alignment = PP_ALIGN.CENTER
        _run_font(sp.runs[0])
    _notes(s, "Expliquer pourquoi le contexte africain appelle une plateforme unifiée.")


def _slide_answer(prs):
    s = slide_new(prs)
    _title(s, "La réponse AfriWonder", size=26)
    col_w, top, gap = 3.15, 1.55, 0.55
    x0 = 0.78
    cols = [
        ("🌐 PWA", "Vite\nReact", 0),
        ("📱 Mobile", "Expo\nReact Native", 1),
        ("🖥 Backend", "Node\nExpress", 2),
    ]
    for label, stack, idx in cols:
        lx = x0 + idx * (col_w + gap)
        _box(s, lx, top, col_w, 2.25, f"{label}\n{stack}", accent=CYAN if idx == 0 else (AMBER if idx == 1 else GREEN), title_size=13)
        if idx < 2:
            _arrow_h(s, lx + col_w + 0.02, top + 1.12, lx + col_w + gap - 0.02)

    cx, cy, r = 5.95, 4.65, 1.18
    hub = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(2 * r), Inches(2 * r))
    hub.fill.solid()
    hub.fill.fore_color.rgb = BG_CARD
    hub.line.color.rgb = CYAN
    hub.line.width = Pt(1)
    hf = hub.text_frame
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = hf.paragraphs[0]
    hp.text = "API unifiée\nExpress + Prisma"
    hp.font.size = Pt(12)
    hp.font.bold = True
    hp.font.color.rgb = TEXT
    _run_font(hp.runs[0])
    hp.alignment = PP_ALIGN.CENTER

    for idx in range(3):
        lx = x0 + idx * (col_w + gap) + col_w / 2
        ly = top + 2.25
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(lx), Inches(ly), Inches(6.58), Inches(cy + r))
        c.line.color.rgb = RGBColor(71, 85, 105)

    _notes(s, "Présenter la triple distribution : web PWA, mobile native et l’API commune.")


def _slide_journey(prs):
    s = slide_new(prs)
    _title(s, "Parcours utilisateur unifié", size=26)
    steps = [
        ("🔑", "Inscription\n(JWT)"),
        ("📷", "Publier une vidéo"),
        ("🛒", "Vendre sur\nMarketplace"),
        ("💼", "Retirer gains\n(Wallet)"),
    ]
    x0, w, top = 0.65, 2.82, 2.05
    for i, (icon, lab) in enumerate(steps):
        left = x0 + i * (w + 0.38)
        ic = s.shapes.add_textbox(Inches(left + w / 2 - 0.35), Inches(top), Inches(0.7), Inches(0.45)).text_frame
        ip = ic.paragraphs[0]
        ip.text = icon
        ip.font.size = Pt(26)
        ip.alignment = PP_ALIGN.CENTER
        _box(s, left, top + 0.55, w, 1.25, lab, accent=CYAN if i % 2 == 0 else AMBER, title_size=13)
        if i < len(steps) - 1:
            _arrow_h(s, left + w + 0.02, top + 1.15, left + w + 0.36)
    _notes(s, "Insister sur la continuité : l’utilisateur ne quitte jamais l’écosystème.")


def _slide_arch_global(prs):
    s = slide_new(prs)
    _title(s, "Architecture technique globale", size=24)
    _box(s, 0.72, 1.35, 2.35, 0.72, "Clients\nPWA / Mobile", accent=CYAN, title_size=11)
    _arrow_h(s, 3.1, 1.71, 3.92)
    _box(s, 3.94, 1.35, 2.62, 0.72, "Nginx\nreverse proxy\n🛡", accent=EDGE, title_size=11)

    y_row = 2.35
    _box(s, 0.88, y_row, 2.92, 0.92, "Express\nbackend × 3", accent=CYAN)
    _box(s, 4.26, y_row, 2.35, 0.92, "Socket.IO\n(temps réel)", accent=AMBER)
    _box(s, 7.3, y_row, 2.22, 0.92, "Redis\ncache + rate", accent=GREEN)

    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.25), Inches(2.06), Inches(5.25), Inches(y_row))
    c.line.color.rgb = EDGE
    _box(s, 2.92, 3.45, 4.45, 0.58, "Prisma ORM", accent=EDGE, title_size=12)

    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.15), Inches(y_row + 0.92), Inches(5.15), Inches(3.45))
    ln.line.color.rgb = EDGE

    _box(s, 2.92, 4.28, 4.55, 0.78, "PostgreSQL", accent=CYAN, title_size=14)

    l2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.18), Inches(4.03), Inches(5.18), Inches(4.28))
    l2.line.color.rgb = EDGE

    _notes(s, "Décrire le routage HTTPS, la scalabilité horizontale des backends et la persistance.")


def _slide_stack_table(prs):
    s = slide_new(prs)
    _title(s, "Stack technique", size=26)
    headers = ("Web (PWA)", "Mobile (Expo)", "Backend")
    cols = ["Vite 6", "React 18", "Tailwind CSS", "TanStack Query", "PWA (Workbox)"]
    colb = ["Expo 54", "RN 0.81", "Zustand", "Agora (lives)", "Expo Router"]
    colc = ["Express", "Prisma", "Zod", "Socket.IO", "JWT"]

    cw = [4.08, 4.08, 4.08]
    x0, y0 = 0.72, 1.55
    for j, h in enumerate(headers):
        _box(s, x0 + sum(cw[:j]), y0, cw[j], 0.52, h, accent=CYAN, title_size=12)
    max_r = max(len(cols), len(colb), len(colc))
    for r in range(max_r):
        y = y0 + 0.54 + r * 0.56
        row_vals = [cols[r] if r < len(cols) else "", colb[r] if r < len(colb) else "", colc[r] if r < len(colc) else ""]
        for j, val in enumerate(row_vals):
            if val:
                _box(
                    s,
                    x0 + sum(cw[:j]),
                    y,
                    cw[j],
                    0.5,
                    val,
                    accent=EDGE,
                    bold=False,
                    title_size=10.5,
                    fill=BG_CARD,
                )
    _notes(s, "Expliquer ces choix : rapidité, langage partagé, écosystème et alignement produit.")


def _slide_data_star(prs):
    s = slide_new(prs)
    _title(s, "Modèle de données (aperçu)", size=26)
    hub_l, hub_t = 5.55, 2.85
    hub = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(hub_l), Inches(hub_t), Inches(2.2), Inches(1.02))
    hub.fill.solid()
    hub.fill.fore_color.rgb = BG_CARD
    hub.line.color.rgb = CYAN
    hub.line.width = Pt(1)
    hf = hub.text_frame
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = hf.paragraphs[0]
    hp.text = "User"
    hp.font.size = Pt(16)
    hp.font.bold = True
    hp.font.color.rgb = TEXT
    _run_font(hp.runs[0])
    hp.alignment = PP_ALIGN.CENTER
    hx = hub_l + 1.1
    hy = hub_t + 0.51
    sats = [
        ("Video", 0.85, 1.45),
        ("Order", 9.65, 1.45),
        ("Wallet", 0.78, 3.85),
        ("Live", 9.88, 3.85),
        ("Message", 0.82, 5.45),
        ("Story", 9.95, 5.45),
    ]
    for lab, lx, ly in sats:
        _box(s, lx, ly, 2.32, 0.58, lab, accent=AMBER, title_size=12)
        sx, sy = lx + 1.16, ly + 0.29
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(sx), Inches(sy), Inches(hx), Inches(hy))
        ln.line.color.rgb = RGBColor(71, 85, 105)

    disc = s.shapes.add_textbox(Inches(0.72), Inches(6.55), Inches(12), Inches(0.55)).text_frame.paragraphs[0]
    disc.text = "+ de 100 modèles (ordre grandeur) · Prisma 7.8"
    disc.font.size = Pt(11)
    disc.font.color.rgb = TEXT_MUTED
    disc.alignment = PP_ALIGN.CENTER
    _run_font(disc.runs[0])
    _notes(s, "Montrer que tout tourne autour de l’utilisateur, avec intégrité transactionnelle où il le faut.")


def _slide_security(prs):
    s = slide_new(prs)
    _title(s, "Sécurité", size=28)
    _box(s, 3.92, 1.52, 5.52, 0.52, "SÉCURITÉ", accent=CYAN, title_size=16)
    rows = [
        (("JWT+jti", CYAN), ("Helmet", AMBER)),
        (("Rate-limit", AMBER), ("CORS", CYAN)),
        (("CSRF", CYAN), ("Proxy · liste blanche\n(anti-SSRF)", AMBER)),
        (("Sanitization", GREEN), ("Sharp · uploads", GREEN)),
    ]
    y = 2.18
    for (ll, al), (lr, ar) in rows:
        _box(s, 1.52, y, 5.05, 0.55, ll, accent=al, bold=True, title_size=11)
        _box(s, 7.0, y, 5.05, 0.55, lr, accent=ar, bold=True, title_size=11)
        y += 0.62
    _notes(s, "Lister vite les briques disponibles dans le projet et insister sur le proxy média avec liste blanche.")


def _slide_realtime(prs):
    s = slide_new(prs)
    _title(s, "Temps réel", size=28)
    _box(s, 1.08, 1.78, 1.92, 1.55, "📱", accent=EDGE, title_size=36)
    _box(s, 9.12, 1.78, 1.92, 1.55, "📱", accent=EDGE, title_size=36)
    for x in (3.52, 4.92, 6.42):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.42), Inches(0.55), Inches(0.4))
        o.fill.solid()
        o.fill.fore_color.rgb = CYAN
        o.line.fill.background()
    _box(s, 4.55, 2.18, 4.28, 0.92, "Socket.IO + Redis\n(rooms / adapter)", accent=CYAN, title_size=12)
    _box(s, 3.35, 3.85, 6.65, 0.95, "Lives : Agora (flux vidéo)", accent=AMBER, title_size=14)
    cap = s.shapes.add_textbox(Inches(0.72), Inches(5.15), Inches(12), Inches(0.9)).text_frame.paragraphs[0]
    cap.text = "Messagerie (Socket.IO) · Lives (Agora)"
    cap.font.size = Pt(14)
    cap.font.color.rgb = TEXT
    cap.alignment = PP_ALIGN.CENTER
    _run_font(cap.runs[0])
    _notes(
        s,
        "Préciser : le chat transite par le backend ; le flux vidéo live est géré côté Agora (P2P / SFU selon config).",
    )


def _slide_devops(prs):
    s = slide_new(prs)
    _title(s, "DevOps & CI/CD", size=26)
    chain = [
        ("Git push", 0.72),
        ("GitHub Actions\nlint · test · build", 2.95),
        ("Docker Compose\nprod", 5.82),
        ("Déploiement", 8.65),
    ]
    for i, (lab, left) in enumerate(chain):
        _box(s, left, 2.05, 2.05, 1.12, lab, accent=CYAN if "Git" in lab else EDGE, title_size=11)
        if i < len(chain) - 1:
            nxt = chain[i + 1][1]
            _arrow_h(s, left + 2.08, 2.61, nxt - 0.08)

    _box(s, 1.22, 3.55, 2.55, 0.68, "Nginx", accent=CYAN)
    _box(s, 4.15, 3.55, 2.35, 0.68, "3× Back", accent=AMBER)
    _box(s, 6.85, 3.55, 2.25, 0.68, "Postgres", accent=GREEN)
    _box(s, 9.25, 3.55, 2.15, 0.68, "Redis", accent=EDGE)

    hint = s.shapes.add_textbox(Inches(0.72), Inches(4.55), Inches(12), Inches(0.55)).text_frame.paragraphs[0]
    hint.text = "CI : PR ≤ 400 lignes · quality gates (cf. ci.yml)"
    hint.font.size = Pt(11)
    hint.font.color.rgb = TEXT_MUTED
    hint.alignment = PP_ALIGN.CENTER
    _run_font(hint.runs[0])

    _notes(s, "Évoquer la CI (contraintes PR, tests, audits) et le déploiement conteneurisé.")


def _slide_ia(prs):
    s = slide_new(prs)
    _title(s, "IA & perspectives", size=26)
    _box(
        s,
        0.78,
        1.65,
        5.78,
        3.55,
        "Aujourd’hui\n• Traduction (LibreTranslate)\n• Chatbot (données)\n• STT / voix — placeholder selon périmètre",
        accent=CYAN,
        title_size=14,
        center=False,
        align_left=True,
    )

    _box(
        s,
        6.88,
        1.65,
        5.68,
        3.55,
        "Demain\n• Modularisation (bounded contexts)\n• STT finalisé · files médias\n• Feature flags / rollout",
        accent=AMBER,
        title_size=14,
        center=False,
        align_left=True,
    )
    _notes(s, "Ne pas survendre l’IA : bilan factuel aujourd’hui, roadmap crédible demain.")


def _slide_difficulties(prs):
    s = slide_new(prs)
    _title(s, "Difficultés & compétences", size=26)
    cards = [
        ("Buffering vidéo\n& réseau instable", CYAN),
        ("Schéma Prisma\nà grande échelle", AMBER),
        ("Cohérence\nPWA / Mobile", GREEN),
    ]
    x0, w = 0.85, 3.72
    for i, (t, ac) in enumerate(cards):
        _box(s, x0 + i * (w + 0.42), 1.92, w, 1.65, t, accent=ac, title_size=13)

    skills = ["Architecture full-stack", "Paiements", "Sécurité", "CI/CD", "Temps réel"]
    gx, gy = 0.75, 4.08
    for i, sk in enumerate(skills):
        row, col = divmod(i, 3)
        _box(s, gx + col * 3.95, gy + row * 0.72, 3.85, 0.58, sk, accent=EDGE, bold=False, title_size=11)
    _notes(s, "Relier chaque difficulté aux apprentissages concrets (réseau, ORM, double client).")


def _slide_conclusion(prs):
    s = slide_new(prs)
    tb = s.shapes.add_textbox(Inches(0.72), Inches(2.15), Inches(12), Inches(2.2)).text_frame
    p0 = tb.paragraphs[0]
    p0.text = "Merci de votre attention"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = TEXT
    _run_font(p0.runs[0])
    p0.alignment = PP_ALIGN.CENTER
    p1 = tb.add_paragraph()
    p1.text = "Des questions ?"
    p1.font.size = Pt(26)
    p1.font.color.rgb = CYAN
    _run_font(p1.runs[0])
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(16)
    p2 = tb.add_paragraph()
    p2.text = "AfriWonder — Une plateforme, toutes les opportunités"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    _run_font(p2.runs[0])
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(22)

    ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(5.15), Inches(2.65), Inches(1.35))
    ph.fill.solid()
    ph.fill.fore_color.rgb = BG_CARD
    ph.line.color.rgb = CYAN
    ph.line.width = Pt(1)
    ph.adjustments[0] = 0.08
    tf = ph.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]
    pp.text = "QR code\n(à insérer)\n→ dépôt / doc"
    pp.font.size = Pt(11)
    pp.font.color.rgb = TEXT_MUTED
    _run_font(pp.runs[0])
    pp.alignment = PP_ALIGN.CENTER

    _notes(
        s,
        "Clôturer calmement. Si autorisé : remplacer le cadre par un vrai QR (Insertion > Image ou complément).",
    )


def build_public_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_visual_deck_core(prs)
    _slide_conclusion(prs)
    return prs


def build_prep_jury_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_visual_deck_core(prs)
    face_n = [0]
    _slide_prep_part_intro(prs, face_n)
    n_sections, n_slides = append_markdown_prep_notes_only(prs, face_n)
    print(f"INFO prep jury : sections MD ~ {n_sections} ; diapos neutres ~ {n_slides}")
    _slide_conclusion(prs)
    return prs


def main() -> None:
    public_only = "--public-only" in sys.argv
    p_pub = build_public_presentation()
    p_pub.save(OUT_PUBLIC)
    print(f"OK PUBLIQUE (a projeter) : {OUT_PUBLIC}")
    if public_only:
        print("(PREP non regenere : relance sans --public-only)")
        return
    p_prep = build_prep_jury_presentation()
    p_prep.save(OUT_PREP)
    print(f"OK PREP personnelle / jury (notes) : {OUT_PREP}")


if __name__ == "__main__":
    main()
